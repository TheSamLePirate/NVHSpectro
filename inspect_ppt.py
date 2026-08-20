import collections
import collections.abc
import pptx
import sys

def analyze_ppt(file_path):
    try:
        prs = pptx.Presentation(file_path)
    except Exception as e:
        print(f"Error opening presentation: {e}")
        return

    print(f"Number of slides: {len(prs.slides)}")
    
    chart_count = 0
    image_count = 0
    
    for i, slide in enumerate(prs.slides):
        print(f"--- Slide {i+1} ---")
        for j, shape in enumerate(slide.shapes):
            if shape.has_chart:
                chart_count += 1
                chart = shape.chart
                print(f"  Shape {j+1}: Chart ({chart.chart_type})")
            elif shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.PICTURE:
                image_count += 1
                print(f"  Shape {j+1}: Image")
            elif shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.GROUP:
                print(f"  Shape {j+1}: Group")
            else:
                print(f"  Shape {j+1}: Other type ({shape.shape_type})")

    print(f"Total charts: {chart_count}")
    print(f"Total images: {image_count}")

if __name__ == '__main__':
    analyze_ppt(sys.argv[1])
