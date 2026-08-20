import pptx
import os
import sys

def extract_images(file_path, output_dir):
    try:
        prs = pptx.Presentation(file_path)
    except Exception as e:
        print(f"Error opening presentation: {e}")
        return

    os.makedirs(output_dir, exist_ok=True)
    
    img_count = 0
    for i, slide in enumerate(prs.slides):
        for j, shape in enumerate(slide.shapes):
            if shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.PICTURE:
                image = shape.image
                image_bytes = image.blob
                image_ext = image.ext
                image_filename = f"slide_{i+1}_img_{j+1}.{image_ext}"
                with open(os.path.join(output_dir, image_filename), 'wb') as f:
                    f.write(image_bytes)
                img_count += 1
                if img_count == 1:
                    print(f"Extracted first image: {image_filename}")

    print(f"Total images extracted: {img_count}")

if __name__ == '__main__':
    extract_images(sys.argv[1], sys.argv[2])
